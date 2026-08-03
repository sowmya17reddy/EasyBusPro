# -*- coding: utf-8 -*-
"""
Builds EasyBus_Pro_Internship_Presentation.pptx
20-slide internship review / project demonstration deck.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn
import copy

# ---------------------------------------------------------------- palette --
DARK_BLUE   = RGBColor(0x0B, 0x1F, 0x3F)   # primary dark navy
DARK_BLUE_2 = RGBColor(0x11, 0x2B, 0x54)   # slightly lighter navy for panels
SKY_BLUE    = RGBColor(0x33, 0xA7, 0xE0)   # accent 1
ORANGE      = RGBColor(0xF2, 0x8C, 0x28)   # accent 2
LIGHT_GRAY  = RGBColor(0xF3, 0xF5, 0xF8)   # slide background
MID_GRAY    = RGBColor(0xDD, 0xE2, 0xE9)   # card borders
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_DARK   = RGBColor(0x20, 0x28, 0x38)
TEXT_MUTED  = RGBColor(0x5A, 0x64, 0x74)

FONT = "Segoe UI"
FONT_SEMIBOLD = "Segoe UI Semibold"
ICON_FONT = "Segoe UI Emoji"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


# ------------------------------------------------------------- utilities --
def I(v):
    """Force any coordinate/dimension to a plain int (EMU). python-pptx does not
    coerce floats when writing XML, and a float such as x='6095847.5' or a float
    cx/cy on a connector is silently accepted by python-pptx/lxml but rejected by
    PowerPoint's own strict OOXML parser ('PowerPoint could not open the file'),
    even though the file re-opens fine in python-pptx. Every value that reaches a
    shape/connector/textbox call must be an int."""
    return int(round(v))


def add_slide():
    return prs.slides.add_slide(BLANK)


def set_background(slide, color):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, color, line_color=None, line_w=None, shadow=False, radius=None):
    x, y, w, h = I(x), I(y), I(w), I(h)
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, x, y, w, h)
    if radius is not None:
        try:
            shp.adjustments[0] = radius
        except Exception:
            pass
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
        shp.line.width = line_w or Pt(0.75)
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, w, h, text, size=18, color=TEXT_DARK, bold=False,
             align=PP_ALIGN.LEFT, font=FONT, anchor=MSO_ANCHOR.TOP, italic=False,
             line_spacing=None, wrap=True):
    x, y, w, h = I(x), I(y), I(w), I(h)
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_spacing:
            p.line_spacing = line_spacing
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.name = font
        r.font.color.rgb = color
    return tb


def add_bullets(slide, x, y, w, h, items, size=15, color=TEXT_DARK, font=FONT,
                 space_after=10, bullet_color=SKY_BLUE, bold_lead=True):
    """items: list of (lead, rest) or plain strings. Draws a small square bullet."""
    x, y, w, h = I(x), I(y), I(w), I(h)
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(space_after)
        p.line_spacing = 1.12
        if isinstance(item, tuple):
            lead, rest = item
            r1 = p.add_run()
            r1.text = "\u25A0  " + lead
            r1.font.size = Pt(size)
            r1.font.bold = True
            r1.font.name = font
            r1.font.color.rgb = TEXT_DARK
            if rest:
                r2 = p.add_run()
                r2.text = "  " + rest
                r2.font.size = Pt(size)
                r2.font.bold = False
                r2.font.name = font
                r2.font.color.rgb = TEXT_MUTED
        else:
            r1 = p.add_run()
            r1.text = "\u25A0  " + item
            r1.font.size = Pt(size)
            r1.font.bold = False
            r1.font.name = font
            r1.font.color.rgb = color
        # color the bullet square only (first 2 chars) - approximate by leaving whole run; acceptable
    return tb


def icon_badge(slide, cx, cy, d, icon, bg=SKY_BLUE, icon_color=WHITE, icon_size=20):
    """cx,cy = center in EMU, d = diameter"""
    cx, cy, d = I(cx), I(cy), I(d)
    x = cx - d // 2
    y = cy - d // 2
    circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, d, d)
    circ.fill.solid()
    circ.fill.fore_color.rgb = bg
    circ.line.fill.background()
    circ.shadow.inherit = False
    tb = slide.shapes.add_textbox(x, y, d, d)
    tf = tb.text_frame
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = icon
    r.font.size = Pt(icon_size)
    r.font.name = ICON_FONT
    r.font.color.rgb = icon_color
    return circ


def header_bar(slide, kicker, title, accent=ORANGE):
    """Standard content-slide header: dark navy bar, kicker + title, orange rule."""
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.35), DARK_BLUE)
    add_text(slide, Inches(0.55), Inches(0.18), Inches(11), Inches(0.32),
              kicker.upper(), size=12, color=SKY_BLUE, bold=True, font=FONT)
    add_text(slide, Inches(0.53), Inches(0.5), Inches(11.5), Inches(0.55),
              title, size=26, color=WHITE, bold=True, font=FONT)
    add_rect(slide, Inches(0.55), Inches(1.14), Inches(0.55), Pt(4), accent)
    return


def footer(slide, num):
    add_text(slide, Inches(12.55), Inches(7.14), Inches(0.6), Inches(0.3),
              str(num), size=11, color=TEXT_MUTED, align=PP_ALIGN.RIGHT)
    add_text(slide, Inches(0.55), Inches(7.14), Inches(4), Inches(0.3),
              "EasyBus Pro  |  Internship Project Review", size=11, color=TEXT_MUTED)


def add_notes(slide, text):
    notes = slide.notes_slide
    notes.notes_text_frame.text = text


def connector_arrow(slide, x1, y1, x2, y2, color=SKY_BLUE, weight=2.25):
    x1, y1, x2, y2 = I(x1), I(y1), I(x2), I(y2)
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    conn.line.color.rgb = color
    conn.line.width = Pt(weight)
    line = conn.line
    ln = line._get_or_add_ln()
    tail = ln.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'})
    ln.append(tail)
    return conn


# =============================================================================
# SLIDE 1 — TITLE
# =============================================================================
s = add_slide()
set_background(s, DARK_BLUE)
add_rect(s, 0, 0, Inches(4.6), SLIDE_H, DARK_BLUE_2)
add_rect(s, Inches(4.6), 0, Pt(3), SLIDE_H, ORANGE)
# decorative circles
c1 = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-1.2), Inches(-1.3), Inches(3.6), Inches(3.6))
c1.fill.solid(); c1.fill.fore_color.rgb = SKY_BLUE; c1.fill.transparency = 0
c1.line.fill.background(); c1.shadow.inherit = False
sp = c1.fill.fore_color._xFill
c1.fill.fore_color.rgb = SKY_BLUE
c2 = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(3.6), Inches(5.6), Inches(2.2), Inches(2.2))
c2.fill.solid(); c2.fill.fore_color.rgb = ORANGE
c2.line.fill.background(); c2.shadow.inherit = False

# bus illustration placeholder panel
add_rect(s, Inches(0.55), Inches(0.55), Inches(3.5), Inches(2.2), DARK_BLUE, line_color=SKY_BLUE, line_w=Pt(1.25), radius=0.08)
add_text(s, Inches(0.75), Inches(1.15), Inches(3.1), Inches(1.1),
          "[ Suggested visual: flat-style bus / road\nillustration or EasyBus Pro logo mark ]",
          size=11.5, color=SKY_BLUE, italic=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

icon_badge(s, Inches(5.6), Inches(1.35), Inches(0.75), "\U0001F68C", bg=ORANGE, icon_size=26)
add_text(s, Inches(5.15), Inches(2.65), Inches(7.6), Inches(1.3), "EasyBus Pro",
          size=54, color=WHITE, bold=True, font=FONT)
add_text(s, Inches(5.18), Inches(3.55), Inches(7.5), Inches(0.6),
          "Bus Ticket Booking System", size=22, color=SKY_BLUE, bold=False)
add_rect(s, Inches(5.2), Inches(4.25), Inches(1.6), Pt(3), ORANGE)

add_text(s, Inches(5.2), Inches(4.55), Inches(7.4), Inches(0.4),
          "Internship Project Presentation", size=15, color=WHITE, bold=True)

details = [
    ("Student", "S. Sowmya  |  2023003449"),
    ("College", "GITAM (Deemed to be University), Hyderabad"),
    ("Department", "Computer Science and Engineering"),
    ("Company", "DayLearner Private Limited, Kakinada"),
    ("Mentor", "<< Mentor Name >>"),
    ("Academic Year", "2025 - 2026"),
]
y = Inches(5.05)
for label, val in details:
    add_text(s, Inches(5.2), y, Inches(2.1), Inches(0.32), label, size=12.5, color=SKY_BLUE, bold=True)
    add_text(s, Inches(7.15), y, Inches(5.4), Inches(0.32), val, size=12.5, color=WHITE)
    y += Inches(0.36)

add_notes(s, (
    "PRESENTER NOTES: Greet the panel, introduce yourself, and state the project name and one-line "
    "purpose: 'EasyBus Pro is a bus ticket booking web application built during my internship at "
    "DayLearner Private Limited using Python, Flask and SQLite3.' Keep this slide on screen for under "
    "30 seconds.\n\n"
    "ANIMATION SUGGESTION: Fade-in the title text, then a subtle wipe-in for the detail rows (top to "
    "bottom, 0.2s stagger).\n"
    "TRANSITION SUGGESTION: Morph or Fade into Slide 2."
))

# =============================================================================
# SLIDE 2 — ABOUT THE INTERNSHIP (timeline)
# =============================================================================
s = add_slide()
set_background(s, LIGHT_GRAY)
header_bar(s, "Internship Overview", "About the Internship")

info = [
    ("\U0001F3E2", "Organization", "DayLearner Private Limited, Kakinada"),
    ("\u23F1", "Duration", "8 weeks (2 months) — 25 May to 25 Jul 2026"),
    ("\U0001F4BB", "Domain", "Python Full-Stack Web Development with AI & Prompt Engineering"),
    ("\U0001F4CB", "Project Assigned", "Design & build a working bus ticket booking web app"),
]
x = Inches(0.55); y = Inches(1.55); card_w = Inches(2.95); gap = Inches(0.18)
for icon, label, val in info:
    add_rect(s, x, y, card_w, Inches(1.55), WHITE, line_color=MID_GRAY, line_w=Pt(1), radius=0.07)
    icon_badge(s, x + Inches(0.5), y + Inches(0.45), Inches(0.62), icon, bg=DARK_BLUE, icon_size=20)
    add_text(s, x + Inches(0.18), y + Inches(0.85), card_w - Inches(0.36), Inches(0.3),
              label.upper(), size=11, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.15), y + Inches(1.13), card_w - Inches(0.3), Inches(0.45),
              val, size=10.5, color=TEXT_DARK, align=PP_ALIGN.CENTER)
    x += card_w + gap

# timeline
add_text(s, Inches(0.55), Inches(3.55), Inches(6), Inches(0.4), "Responsibilities Timeline", size=16, bold=True, color=DARK_BLUE)
tl_y = Inches(4.35)
add_rect(s, Inches(0.75), tl_y, Inches(11.8), Pt(3), MID_GRAY)
weeks = [
    ("W1-2", "Orientation, Python/Flask setup,\ndatabase design"),
    ("W3-4", "Landing page, registration,\nlogin, sessions"),
    ("W5-6", "Search, results, seat\nselection, payment UI"),
    ("W7-8", "Booking confirmation, testing,\ndocumentation, demo"),
]
seg = Inches(11.8) / 4
for i, (wk, desc) in enumerate(weeks):
    cx = I(Inches(0.75) + seg * i + seg / 2)
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL, I(cx - Inches(0.11)), I(tl_y - Inches(0.08)), Inches(0.22), Inches(0.22))
    dot.fill.solid(); dot.fill.fore_color.rgb = SKY_BLUE if i % 2 == 0 else ORANGE
    dot.line.color.rgb = WHITE; dot.line.width = Pt(2); dot.shadow.inherit = False
    add_text(s, cx - Inches(1.35), tl_y + Inches(0.25), Inches(2.7), Inches(0.35), wk,
              size=13, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)
    add_text(s, cx - Inches(1.35), tl_y + Inches(0.62), Inches(2.7), Inches(0.9), desc,
              size=10.5, color=TEXT_MUTED, align=PP_ALIGN.CENTER)

footer(s, 2)
add_notes(s, (
    "PRESENTER NOTES: Briefly explain the internship structure — 8 weeks at DayLearner Private "
    "Limited, working under weekly milestone reviews. Mention that the objective was to learn the "
    "full web development cycle (not just theory) by building one complete project end-to-end.\n\n"
    "ANIMATION SUGGESTION: Animate the four info cards with a 'Fly in from bottom' effect (staggered "
    "0.15s each), then animate the timeline dots left-to-right with 'Appear'.\n"
    "TRANSITION SUGGESTION: Push (left) from Slide 1."
))

# =============================================================================
# SLIDE 3 — PROJECT INTRODUCTION (infographic)
# =============================================================================
s = add_slide()
set_background(s, LIGHT_GRAY)
header_bar(s, "Project Introduction", "Why EasyBus Pro?")

add_rect(s, Inches(0.55), Inches(1.5), Inches(6.0), Inches(5.35), WHITE, line_color=MID_GRAY, radius=0.05)
add_text(s, Inches(0.85), Inches(1.75), Inches(5.4), Inches(0.4), "Problem Statement", size=16, bold=True, color=DARK_BLUE)
add_text(s, Inches(0.85), Inches(2.2), Inches(5.4), Inches(1.9),
          "Most small & medium bus operators still rely on counter or phone-based booking. "
          "Passengers cannot compare operators, cannot see seat layouts, and get only a paper "
          "slip as proof — with no digital record of their journey.",
          size=13.5, color=TEXT_MUTED, line_spacing=1.25)

add_text(s, Inches(0.85), Inches(4.15), Inches(5.4), Inches(0.4), "Objective", size=16, bold=True, color=DARK_BLUE)
add_bullets(s, Inches(0.85), Inches(4.6), Inches(5.4), Inches(2.1), [
    "Build a working, database-driven booking web app",
    "Cover the full cycle: search \u2192 select \u2192 pay \u2192 confirm",
    "Apply Flask, SQLite and web fundamentals practically",
], size=13.5)

# right infographic stat blocks
stats = [
    ("\U0001F3AB", "1", "Single platform to search, compare & book"),
    ("\U0001F5C3", "3", "Linked database tables (users, bus, bookings)"),
    ("\u23F1", "8", "Weeks to design, build, test & document"),
]
x = Inches(6.85); y = Inches(1.5)
for icon, num, desc in stats:
    add_rect(s, x, y, Inches(5.9), Inches(1.65), DARK_BLUE, radius=0.09)
    icon_badge(s, x + Inches(0.85), y + Inches(0.82), Inches(1.0), icon, bg=ORANGE, icon_size=30)
    add_text(s, x + Inches(1.55), y + Inches(0.2), Inches(1.0), Inches(0.8), num,
              size=36, bold=True, color=SKY_BLUE)
    add_text(s, x + Inches(2.7), y + Inches(0.32), Inches(3.0), Inches(1.0), desc,
              size=13, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    y += Inches(1.85)

footer(s, 3)
add_notes(s, (
    "PRESENTER NOTES: State the real-world problem in one sentence, then pivot to what the project "
    "achieves. Keep the tone practical — this is the 'why' slide before you show 'how'.\n\n"
    "ANIMATION SUGGESTION: Left card content 'Fade' in as one block; right stat cards 'Fly in from "
    "right' staggered.\n"
    "TRANSITION SUGGESTION: Fade."
))

# =============================================================================
# SLIDE 4 — EXISTING SYSTEM (comparison)
# =============================================================================
s = add_slide()
set_background(s, LIGHT_GRAY)
header_bar(s, "Background", "Existing System & Its Limitations")

add_text(s, Inches(0.55), Inches(1.45), Inches(11), Inches(0.4),
          "Traditional Counter / Telephone Booking Process", size=16, bold=True, color=DARK_BLUE)

steps = ["Passenger\ncalls / visits\ncounter", "Operator checks\nregister for\navailability",
         "Fare & seat\nquoted\nverbally", "Seat noted by\nhand, cash\ncollected", "Paper slip\nissued"]
x = Inches(0.55); seg_w = Inches(2.42)
for i, st in enumerate(steps):
    add_rect(s, x, Inches(2.0), seg_w - Inches(0.15), Inches(1.15), WHITE, line_color=MID_GRAY, radius=0.1)
    add_text(s, x + Inches(0.1), Inches(2.13), seg_w - Inches(0.35), Inches(0.9), st,
              size=11.5, color=TEXT_DARK, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    if i < len(steps) - 1:
        connector_arrow(s, x + seg_w - Inches(0.13), Inches(2.57), x + seg_w + Inches(0.02), Inches(2.57), color=ORANGE)
    x += seg_w

add_text(s, Inches(0.55), Inches(3.55), Inches(11), Inches(0.4), "Key Limitations", size=16, bold=True, color=DARK_BLUE)
probs = [
    ("\u23F0", "Restricted Hours", "Bookings possible only during counter/call hours"),
    ("\U0001F500", "No Comparison", "Each operator must be contacted separately"),
    ("\u274C", "Manual Errors", "Double booking & illegible register entries"),
    ("\U0001F4C4", "No Digital Record", "Paper slip is easily lost, no booking history"),
]
x = Inches(0.55); y = Inches(4.05); w = Inches(2.95)
for icon, title, desc in probs:
    add_rect(s, x, y, w, Inches(2.6), WHITE, line_color=MID_GRAY, radius=0.07)
    icon_badge(s, x + w/2, y + Inches(0.55), Inches(0.7), icon, bg=RGBColor(0xC0, 0x39, 0x2B), icon_size=24)
    add_text(s, x + Inches(0.15), y + Inches(1.05), w - Inches(0.3), Inches(0.4), title,
              size=13.5, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.2), y + Inches(1.45), w - Inches(0.4), Inches(1.0), desc,
              size=11, color=TEXT_MUTED, align=PP_ALIGN.CENTER)
    x += w + Inches(0.15)

footer(s, 4)
add_notes(s, (
    "PRESENTER NOTES: Walk through the manual process quickly (5 steps), then highlight the four "
    "problem cards. Use a real example if asked: 'I had to call three operators separately to compare "
    "fares before a trip.'\n\n"
    "ANIMATION SUGGESTION: Animate the 5-step flow left-to-right with connecting arrows using "
    "'Wipe' entrance; problem cards 'Fly in from bottom' staggered.\n"
    "TRANSITION SUGGESTION: Push (left)."
))

# =============================================================================
# SLIDE 5 — PROPOSED SYSTEM (flow illustration)
# =============================================================================
s = add_slide()
set_background(s, LIGHT_GRAY)
header_bar(s, "Solution", "Proposed System — EasyBus Pro")

add_text(s, Inches(0.55), Inches(1.45), Inches(12.2), Inches(0.6),
          "A single web application that replaces every manual step above with an automated, "
          "database-backed flow.", size=14, color=TEXT_MUTED)

flow = ["Search", "Compare", "Select Seats", "Pay Online", "Confirmed\nE-Ticket"]
x = Inches(0.55); seg_w = Inches(2.42)
for i, st in enumerate(flow):
    add_rect(s, x, Inches(2.25), seg_w - Inches(0.15), Inches(1.1), DARK_BLUE, radius=0.15)
    add_text(s, x + Inches(0.1), Inches(2.4), seg_w - Inches(0.35), Inches(0.8), st,
              size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    if i < len(flow) - 1:
        connector_arrow(s, x + seg_w - Inches(0.13), Inches(2.8), x + seg_w + Inches(0.02), Inches(2.8), color=ORANGE, weight=2.5)
    x += seg_w

add_text(s, Inches(0.55), Inches(3.85), Inches(11), Inches(0.4), "Advantages", size=16, bold=True, color=DARK_BLUE)
adv = [
    ("\u26A1", "Automation", "No manual register updates; data flows straight into the database"),
    ("\U0001F510", "Secure Login", "Session-based authentication protects personal booking data"),
    ("\U0001F5C4", "Database Integration", "Every booking permanently linked to user & bus records"),
    ("\U0001F3A8", "Modern UI", "Clean, responsive interface usable on desktop and mobile"),
]
x = Inches(0.55); y = Inches(4.3); w = Inches(2.95)
for icon, title, desc in adv:
    add_rect(s, x, y, w, Inches(2.55), WHITE, line_color=MID_GRAY, radius=0.07)
    icon_badge(s, x + w/2, y + Inches(0.55), Inches(0.7), icon, bg=SKY_BLUE, icon_size=24)
    add_text(s, x + Inches(0.15), y + Inches(1.05), w - Inches(0.3), Inches(0.4), title,
              size=13.5, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.2), y + Inches(1.45), w - Inches(0.4), Inches(1.0), desc,
              size=10.5, color=TEXT_MUTED, align=PP_ALIGN.CENTER)
    x += w + Inches(0.15)

footer(s, 5)
add_notes(s, (
    "PRESENTER NOTES: Mirror the structure of the previous slide so the panel sees a direct "
    "before/after contrast. Emphasize that every advantage maps to a limitation just shown.\n\n"
    "ANIMATION SUGGESTION: Flow boxes animate with 'Zoom' one after another; advantage cards use "
    "'Fade' staggered.\n"
    "TRANSITION SUGGESTION: Morph (to visually connect from Slide 4's flow into this one)."
))

# =============================================================================
# SLIDE 6 — TECHNOLOGY STACK
# =============================================================================
s = add_slide()
set_background(s, LIGHT_GRAY)
header_bar(s, "Tech Stack", "Technologies Used")

techs = [
    ("\U0001F40D", "Python", "Core programming language — simple syntax, huge ecosystem"),
    ("\U0001F3F4", "Flask", "Lightweight web framework — full control over routes & logic"),
    ("\U0001F5C4", "SQLite3", "Zero-config relational database bundled with Python"),
    ("\U0001F310", "HTML5", "Semantic page structure for every screen"),
    ("\U0001F3A8", "CSS3", "Hand-written responsive styling, no framework dependency"),
    ("\U0001F9E9", "Jinja2", "Server-side templating with reusable base layout"),
    ("\U0001F4BB", "VS Code", "Primary editor — integrated terminal & Python tooling"),
]
cols = 4
cw = Inches(2.95); ch = Inches(2.55); gx = Inches(0.15); gy = Inches(0.2)
x0 = Inches(0.55); y0 = Inches(1.5)
for i, (icon, name, desc) in enumerate(techs):
    col = i % cols
    row = i // cols
    x = x0 + col * (cw + gx)
    y = y0 + row * (ch + gy)
    add_rect(s, x, y, cw, ch, WHITE, line_color=MID_GRAY, radius=0.08)
    icon_badge(s, x + cw/2, y + Inches(0.6), Inches(0.75), icon, bg=DARK_BLUE, icon_size=26)
    add_text(s, x + Inches(0.1), y + Inches(1.15), cw - Inches(0.2), Inches(0.35), name,
              size=15, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.15), y + Inches(1.55), cw - Inches(0.3), Inches(0.9), desc,
              size=10, color=TEXT_MUTED, align=PP_ALIGN.CENTER)

footer(s, 6)
add_notes(s, (
    "PRESENTER NOTES: For each technology, give ONE reason it was chosen (already written on-card). "
    "If asked 'why not Django/Bootstrap?', explain the internship goal was to learn fundamentals "
    "directly (raw SQL, hand-written CSS) rather than rely on an abstraction layer.\n\n"
    "ANIMATION SUGGESTION: Grid cards 'Wipe' in row by row.\n"
    "TRANSITION SUGGESTION: Push (left)."
))

# =============================================================================
# SLIDE 7 — SYSTEM ARCHITECTURE
# =============================================================================
s = add_slide()
set_background(s, LIGHT_GRAY)
header_bar(s, "Architecture", "System Architecture")

layers = [
    ("\U0001F5A5", "USER / BROWSER", "HTML5 + CSS3 pages rendered via Jinja2", SKY_BLUE),
    ("\u2699", "FLASK APPLICATION", "15 routes handle requests, sessions & logic", DARK_BLUE),
    ("\U0001F5C4", "SQLITE DATABASE", "users, bus, bookings tables (database.db)", ORANGE),
    ("\u2705", "BOOKING CONFIRMED", "E-ticket generated & stored permanently", RGBColor(0x2E, 0x9E, 0x5B)),
]
box_w = Inches(6.9); box_h = Inches(1.05)
x = Inches(0.55)
y = Inches(1.55)
for i, (icon, title, desc, color) in enumerate(layers):
    add_rect(s, x, y, box_w, box_h, color, radius=0.12)
    icon_badge(s, x + Inches(0.65), y + box_h/2, Inches(0.6), icon, bg=WHITE, icon_color=color, icon_size=20)
    add_text(s, x + Inches(1.15), y + Inches(0.16), Inches(4.8), Inches(0.35), title,
              size=15, bold=True, color=WHITE)
    add_text(s, x + Inches(1.15), y + Inches(0.55), Inches(5.6), Inches(0.4), desc,
              size=11, color=WHITE)
    if i < len(layers) - 1:
        connector_arrow(s, x + box_w/2, y + box_h + Inches(0.03), x + box_w/2, y + box_h + Inches(0.27), color=DARK_BLUE, weight=2.5)
    y += box_h + Inches(0.4)

add_rect(s, Inches(7.85), Inches(1.55), Inches(4.95), Inches(5.15), WHITE, line_color=MID_GRAY, radius=0.05)
add_text(s, Inches(8.15), Inches(1.78), Inches(4.4), Inches(0.4), "How It Works", size=15, bold=True, color=DARK_BLUE)
add_text(s, Inches(8.15), Inches(2.3), Inches(4.4), Inches(4.2),
          "1.  Browser sends an HTTP\n     request (GET / POST)\n\n"
          "2.  Flask matches the route and\n     runs the view function\n\n"
          "3.  A parameterised SQL query\n     reads or writes SQLite\n\n"
          "4.  Jinja2 renders the final\n     HTML page\n\n"
          "5.  Response is sent back to\n     the browser",
          size=12, color=TEXT_DARK, line_spacing=1.2)

footer(s, 7)
add_notes(s, (
    "PRESENTER NOTES: Explain this as the request-response cycle. Point out that the browser NEVER "
    "talks to the database directly — Flask is the only layer allowed to do that. This is the "
    "three-tier architecture (presentation / application / data).\n\n"
    "ANIMATION SUGGESTION: Reveal each layer box top-to-bottom with 'Fade', arrows appearing right "
    "after each box; right-hand steps appear in sync with each layer.\n"
    "TRANSITION SUGGESTION: Fade."
))

# =============================================================================
# SLIDE 8 — DATABASE DESIGN (ER diagram)
# =============================================================================
s = add_slide()
set_background(s, LIGHT_GRAY)
header_bar(s, "Data Model", "Database Design — ER Diagram")

def table_box(slide, x, y, w, h, title, fields, header_color=DARK_BLUE):
    add_rect(slide, x, y, w, h, WHITE, line_color=MID_GRAY, line_w=Pt(1.25), radius=0.04)
    add_rect(slide, x, y, w, Inches(0.42), header_color, radius=0.0)
    add_text(slide, x, y + Inches(0.02), w, Inches(0.38), title, size=13.5, bold=True,
              color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    fy = y + Inches(0.55)
    for f in fields:
        add_text(slide, x + Inches(0.18), fy, w - Inches(0.3), Inches(0.28), f, size=11,
                  color=TEXT_DARK)
        fy += Inches(0.29)

table_box(s, Inches(0.55), Inches(1.55), Inches(3.15), Inches(2.9), "USERS",
          ["PK  id", "     name", "     email (unique)", "     phone", "     password"], SKY_BLUE)
table_box(s, Inches(5.1), Inches(1.55), Inches(3.15), Inches(3.5), "BOOKINGS",
          ["PK  booking_id", "FK  user_id", "FK  bus_id", "     journey_date", "     seat_numbers",
           "     passengers", "     amount", "     payment_method"], DARK_BLUE)
table_box(s, Inches(9.65), Inches(1.55), Inches(3.15), Inches(3.2), "BUS",
          ["PK  id", "     bus_name", "     from_city / to_city", "     departure / arrival",
           "     price", "     available_seats"], ORANGE)

connector_arrow(s, Inches(3.7), Inches(2.6), Inches(5.1), Inches(2.6), color=TEXT_MUTED, weight=1.75)
add_text(s, Inches(3.72), Inches(2.28), Inches(1.4), Inches(0.3), "1 : N", size=10.5, bold=True, color=TEXT_MUTED)
connector_arrow(s, Inches(9.65), Inches(2.6), Inches(8.25), Inches(2.6), color=TEXT_MUTED, weight=1.75)
add_text(s, Inches(8.35), Inches(2.28), Inches(1.4), Inches(0.3), "N : 1", size=10.5, bold=True, color=TEXT_MUTED)

add_rect(s, Inches(0.55), Inches(5.15), Inches(12.25), Inches(1.55), WHITE, line_color=MID_GRAY, radius=0.06)
add_text(s, Inches(0.85), Inches(5.32), Inches(11.6), Inches(0.35), "Key Design Points", size=14, bold=True, color=DARK_BLUE)
add_bullets(s, Inches(0.85), Inches(5.68), Inches(11.6), Inches(0.95), [
    "Every booking stores bus_id and user_id — no duplicated text data across rows",
    "email carries a UNIQUE constraint, enforced by SQLite itself at insert time",
], size=12.5, space_after=6)

footer(s, 8)
add_notes(s, (
    "PRESENTER NOTES: Point at BOOKINGS in the middle and explain it is the 'linking' table — one "
    "user can have many bookings, one bus can appear in many bookings. This is exactly the 1:N "
    "relationship pattern taught in DBMS.\n\n"
    "ANIMATION SUGGESTION: USERS and BUS tables 'Fly in' from left/right respectively, BOOKINGS "
    "'Fades' in the centre, then both relationship arrows 'Wipe' in.\n"
    "TRANSITION SUGGESTION: Push (left)."
))

# =============================================================================
# SLIDE 9 — PROJECT WORKFLOW
# =============================================================================
s = add_slide()
set_background(s, LIGHT_GRAY)
header_bar(s, "End-to-End Flow", "Project Workflow")

steps = ["Landing\nPage", "Login", "Search\nBus", "Bus\nResults", "Seat\nSelection",
         "Payment", "Booking\nConfirmation", "My\nBookings"]
cols = 4
cw = Inches(2.85); ch = Inches(1.15); gx = Inches(0.25); gy = Inches(0.55)
x0 = Inches(0.75); y0 = Inches(1.6)
positions = []
for i, st in enumerate(steps):
    col = i % cols
    row = i // cols
    x = x0 + col * (cw + gx)
    y = y0 + row * (ch + gy)
    positions.append((x, y))
    add_rect(s, x, y, cw, ch, DARK_BLUE if row == 0 else SKY_BLUE, radius=0.18)
    add_text(s, x, y, cw, ch, st, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    num_badge = s.shapes.add_shape(MSO_SHAPE.OVAL, x - Inches(0.16), y - Inches(0.16), Inches(0.34), Inches(0.34))
    num_badge.fill.solid(); num_badge.fill.fore_color.rgb = ORANGE
    num_badge.line.color.rgb = WHITE; num_badge.line.width = Pt(1.5); num_badge.shadow.inherit = False
    tb = s.shapes.add_textbox(x - Inches(0.16), y - Inches(0.16), Inches(0.34), Inches(0.34))
    tf = tb.text_frame; tf.margin_left=0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = str(i+1); r.font.size = Pt(12); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = FONT

# arrows within row 1 (0->1->2->3) and row 2 (4->5->6->7), plus wrap connector 3->4
for i in range(3):
    x, y = positions[i]
    connector_arrow(s, x + cw + Inches(0.02), y + ch/2, x + cw + gx - Inches(0.02), y + ch/2, color=ORANGE, weight=2.25)
for i in range(4, 7):
    x, y = positions[i]
    connector_arrow(s, x + cw + Inches(0.02), y + ch/2, x + cw + gx - Inches(0.02), y + ch/2, color=ORANGE, weight=2.25)
# short "continues below" stub arrow under the last box of row 1 (step 4)
x3, y3 = positions[3]
connector_arrow(s, x3 + cw/2, y3 + ch + Inches(0.02), x3 + cw/2, y3 + ch + Inches(0.32), color=ORANGE, weight=2.25)

footer(s, 9)
add_notes(s, (
    "PRESENTER NOTES: Walk through the numbered path exactly as a real user would experience it. "
    "Mention that state (login session) and journey data (route/seats/amount) travel forward through "
    "this entire chain until the booking is written to the database.\n\n"
    "ANIMATION SUGGESTION: Reveal boxes 1-8 in numeric order using 'Appear', arrows drawing in "
    "between each ('Wipe').\n"
    "TRANSITION SUGGESTION: Push (left)."
))

# =============================================================================
# SLIDE 10 — PROJECT MODULES (cards)
# =============================================================================
s = add_slide()
set_background(s, LIGHT_GRAY)
header_bar(s, "Modules", "Project Modules at a Glance")

modules = [
    ("\U0001F3E0", "Landing Page", "Hero, search box & feature highlights"),
    ("\U0001F4DD", "Registration", "New account with validation"),
    ("\U0001F511", "Login", "Session-based authentication"),
    ("\U0001F4CA", "Dashboard", "Personalised post-login home"),
    ("\U0001F50D", "Search", "Route, date & passenger input"),
    ("\U0001F4CB", "Results", "Matching buses with fare & timing"),
    ("\U0001F4BA", "Seat Selection", "Interactive coach seat layout"),
    ("\U0001F4B3", "Payment", "UPI / Card / Net Banking choice"),
    ("\u2705", "Booking Confirm.", "E-ticket generated & saved"),
    ("\U0001F4C4", "My Bookings", "Full booking history via JOIN"),
    ("\U0001F464", "Profile", "Account details of user"),
    ("\U0001F6AA", "Logout", "Clears session securely"),
]
cols = 6
cw = Inches(1.98); ch = Inches(1.85); gx = Inches(0.1); gy = Inches(0.15)
x0 = Inches(0.55); y0 = Inches(1.5)
for i, (icon, name, desc) in enumerate(modules):
    col = i % cols
    row = i // cols
    x = x0 + col * (cw + gx)
    y = y0 + row * (ch + gy)
    add_rect(s, x, y, cw, ch, WHITE, line_color=MID_GRAY, radius=0.1)
    icon_badge(s, x + cw/2, y + Inches(0.45), Inches(0.55), icon, bg=DARK_BLUE, icon_size=18)
    add_text(s, x + Inches(0.08), y + Inches(0.82), cw - Inches(0.16), Inches(0.45), name,
              size=10.5, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.1), y + Inches(1.24), cw - Inches(0.2), Inches(0.55), desc,
              size=8.5, color=TEXT_MUTED, align=PP_ALIGN.CENTER)

footer(s, 10)
add_notes(s, (
    "PRESENTER NOTES: Do not read every card aloud — summarise: 'The application has 12 functional "
    "modules covering the complete journey from registration to logout, each backed by its own Flask "
    "route and template.' Point to 2-3 cards only (e.g. Seat Selection, Booking Confirmation).\n\n"
    "ANIMATION SUGGESTION: Grid 'Wipe' in row by row, fast (0.05s stagger) since there are 12 cards.\n"
    "TRANSITION SUGGESTION: Fade."
))

# =============================================================================
# SLIDE 11 — USER INTERFACE (screenshots)
# =============================================================================
s = add_slide()
set_background(s, LIGHT_GRAY)
header_bar(s, "Screens", "User Interface Walkthrough")

import os
SHOT_DIR = r"c:\Users\reddy\OneDrive\Desktop\EasyBusPro\Internship_Report\screenshots"
shots = [
    ("Landing Page", "fig_4_1_home.png", "First screen — search bar & highlights"),
    ("Dashboard", "fig_4_4_dashboard.png", "Post-login personal home screen"),
    ("Search Page", "fig_4_5_search.png", "Route, date & passenger input form"),
    ("Seat Selection", "fig_4_7_seat_selection.png", "Click-to-select interactive coach layout"),
    ("Payment", "fig_4_8_payment.png", "UPI / Card / Net Banking checkout"),
    ("Success Page", "fig_4_9_success.png", "Confirmed e-ticket after booking"),
    ("My Bookings", "fig_4_10_my_bookings.png", "Complete booking history (JOIN query)"),
]
cols = 4
cw = Inches(2.95); ch = Inches(1.55); gx = Inches(0.15)
caption_h = Inches(0.78)   # room reserved below each image for name + one-line description
row_pitch = ch + caption_h + Inches(0.22)
x0 = Inches(0.55); y0 = Inches(1.45)
for i, (name, filename, desc) in enumerate(shots):
    col = i % cols
    row = i // cols
    x = x0 + col * (cw + gx)
    y = y0 + row * row_pitch
    add_rect(s, x, y, cw, ch, WHITE, line_color=MID_GRAY, line_w=Pt(1.25), radius=0.05)
    img_path = os.path.join(SHOT_DIR, filename)
    pad = Inches(0.08)
    frame_w = cw - pad * 2
    frame_h = ch - pad * 2
    if os.path.exists(img_path):
        from PIL import Image as PILImage
        with PILImage.open(img_path) as im:
            iw, ih = im.size
        ratio = min(frame_w / iw, frame_h / ih)
        pic_w = int(iw * ratio)
        pic_h = int(ih * ratio)
        pic_x = int(x + pad + (frame_w - pic_w) / 2)
        pic_y = int(y + pad + (frame_h - pic_h) / 2)
        s.shapes.add_picture(img_path, pic_x, pic_y, width=pic_w, height=pic_h)
    else:
        add_text(s, x + Inches(0.2), y + Inches(0.6), cw - Inches(0.4), Inches(0.6),
                  "[ Screenshot:\n" + name + " ]", size=10, italic=True, color=TEXT_MUTED,
                  align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x + Inches(0.1), y + ch + Inches(0.08), cw - Inches(0.2), Inches(0.3), name,
              size=11.5, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.1), y + ch + Inches(0.4), cw - Inches(0.2), Inches(0.45), desc,
              size=9, color=TEXT_MUTED, align=PP_ALIGN.CENTER)

footer(s, 11)
add_notes(s, (
    "PRESENTER NOTES: These are live screenshots captured directly from the running application "
    "(same set used in the written internship report). Briefly narrate each screen in the order "
    "shown — 15-20 seconds total, since the live demo slide covers this in more depth.\n\n"
    "ANIMATION SUGGESTION: Screenshot cards 'Fade' in row by row.\n"
    "TRANSITION SUGGESTION: Push (left)."
))

# =============================================================================
# SLIDE 12 — DATABASE INTEGRATION
# =============================================================================
s = add_slide()
set_background(s, LIGHT_GRAY)
header_bar(s, "Backend", "Database Integration")

items = [
    ("\U0001F50C", "SQLite Connection", "sqlite3.connect() opens database.db; closed after every request"),
    ("\U0001F504", "CRUD Operations", "Create (INSERT), Read (SELECT + JOIN) used; Update/Delete planned"),
    ("\U0001F510", "Authentication", "Parameterised SELECT checks email & password on login"),
    ("\U0001F3AB", "Booking Storage", "Confirmed bookings INSERTed with user_id & bus_id references"),
    ("\U0001F36A", "Session Handling", "Flask session cookie (signed) tracks the logged-in user"),
]
y = Inches(1.55)
for icon, title, desc in items:
    add_rect(s, Inches(0.55), y, Inches(12.25), Inches(0.92), WHITE, line_color=MID_GRAY, radius=0.12)
    icon_badge(s, Inches(1.15), y + Inches(0.46), Inches(0.6), icon, bg=SKY_BLUE, icon_size=20)
    add_text(s, Inches(1.65), y + Inches(0.13), Inches(3.2), Inches(0.35), title, size=14, bold=True, color=DARK_BLUE)
    add_text(s, Inches(1.65), y + Inches(0.48), Inches(10.8), Inches(0.38), desc, size=11.5, color=TEXT_MUTED)
    y += Inches(1.05)

footer(s, 12)
add_notes(s, (
    "PRESENTER NOTES: Emphasize that all SQL statements use ? placeholders (parameterised queries), "
    "which prevents SQL injection — demonstrated live during testing. Mention that UPDATE/DELETE are "
    "listed under Future Enhancements (profile edit, cancellation).\n\n"
    "ANIMATION SUGGESTION: Rows 'Fly in from left' top to bottom, 0.1s stagger.\n"
    "TRANSITION SUGGESTION: Fade."
))

# =============================================================================
# SLIDE 13 — CHALLENGES FACED (problem/solution)
# =============================================================================
s = add_slide()
set_background(s, LIGHT_GRAY)
header_bar(s, "Real Difficulties", "Challenges Faced & Solutions")

pairs = [
    ("TemplateNotFound error", "Moved all HTML files into the required templates/ folder"),
    ("OperationalError: no such table", "Re-ran database.py after adding the bookings table"),
    ("Duplicate seed data on re-run", "Switched INSERT to INSERT OR IGNORE for idempotency"),
    ("from is a Python reserved word", "Built that particular link inside the Jinja2 template instead of app.py"),
    ("Booking duplicated on page refresh", "Documented as a known issue; Post/Redirect/Get fix planned"),
    ("Wrong seat total (string concat)", "Wrapped values with Number() before arithmetic in JavaScript"),
]
cols = 2
cw = Inches(6.05); ch = Inches(1.55); gx = Inches(0.15); gy = Inches(0.15)
x0 = Inches(0.55); y0 = Inches(1.5)
for i, (prob, sol) in enumerate(pairs):
    col = i % cols
    row = i // cols
    x = x0 + col * (cw + gx)
    y = y0 + row * (ch + gy)
    add_rect(s, x, y, cw, ch, WHITE, line_color=MID_GRAY, radius=0.06)
    icon_badge(s, x + Inches(0.45), y + Inches(0.4), Inches(0.5), "\u26A0", bg=RGBColor(0xC0,0x39,0x2B), icon_size=16)
    add_text(s, x + Inches(0.8), y + Inches(0.12), cw - Inches(1.0), Inches(0.55), prob,
              size=12, bold=True, color=DARK_BLUE)
    icon_badge(s, x + Inches(0.45), y + Inches(1.12), Inches(0.5), "\u2705", bg=RGBColor(0x2E,0x9E,0x5B), icon_size=16)
    add_text(s, x + Inches(0.8), y + Inches(0.85), cw - Inches(1.0), Inches(0.6), sol,
              size=10.5, color=TEXT_MUTED)

footer(s, 13)
add_notes(s, (
    "PRESENTER NOTES: Pick 2-3 of these to narrate in detail if time-limited (the 'from' keyword issue "
    "is a good technical talking point — shows genuine debugging depth). State that every fix came "
    "from reading the actual error message, not guessing.\n\n"
    "ANIMATION SUGGESTION: Cards 'Fade' in two-by-two (row by row).\n"
    "TRANSITION SUGGESTION: Push (left)."
))

# =============================================================================
# SLIDE 14 — TESTING (checklist)
# =============================================================================
s = add_slide()
set_background(s, LIGHT_GRAY)
header_bar(s, "Quality Assurance", "Testing")

tests = [
    ("Registration Testing", "Duplicate e-mail rejected; password-mismatch rejected"),
    ("Login Testing", "Valid/invalid credentials; SQL-injection attempt blocked"),
    ("Search Testing", "Valid route returns buses; unknown route shows clear message"),
    ("Booking Testing", "Full flow (search \u2192 seats \u2192 pay \u2192 confirm) verified end-to-end"),
    ("Database Testing", "Verified rows in database.db using DB Browser for SQLite"),
]
y = Inches(1.6)
for title, desc in tests:
    icon_badge(s, Inches(1.0), y + Inches(0.35), Inches(0.55), "\u2714", bg=RGBColor(0x2E,0x9E,0x5B), icon_size=20)
    add_rect(s, Inches(1.45), y, Inches(11.3), Inches(0.85), WHITE, line_color=MID_GRAY, radius=0.14)
    add_text(s, Inches(1.75), y + Inches(0.1), Inches(3.5), Inches(0.35), title, size=14, bold=True, color=DARK_BLUE)
    add_text(s, Inches(1.75), y + Inches(0.46), Inches(10.8), Inches(0.35), desc, size=11, color=TEXT_MUTED)
    y += Inches(1.0)

footer(s, 14)
add_notes(s, (
    "PRESENTER NOTES: State that testing was manual (unit + functional + integration + user "
    "acceptance) and reference the test-case tables in the written report (Tables 4.13-4.15). Mention "
    "automated testing with pytest as a noted gap / future improvement.\n\n"
    "ANIMATION SUGGESTION: Checkmarks 'Appear' one by one top to bottom with a short pause, like a "
    "checklist being ticked off live.\n"
    "TRANSITION SUGGESTION: Fade."
))

# =============================================================================
# SLIDE 15 — LEARNING OUTCOMES
# =============================================================================
s = add_slide()
set_background(s, LIGHT_GRAY)
header_bar(s, "Growth", "Learning Outcomes")

add_rect(s, Inches(0.55), Inches(1.5), Inches(6.0), Inches(5.35), WHITE, line_color=MID_GRAY, radius=0.05)
add_text(s, Inches(0.85), Inches(1.72), Inches(5.4), Inches(0.4), "Technical Skills", size=16, bold=True, color=DARK_BLUE)
add_bullets(s, Inches(0.85), Inches(2.2), Inches(5.4), Inches(4.3), [
    "Flask routing, sessions & the request-response cycle",
    "SQLite database design, CRUD & JOIN queries",
    "HTML5, hand-written CSS3 (Flexbox / Grid)",
    "Jinja2 template inheritance",
    "Systematic debugging using tracebacks",
], size=13.5, space_after=14)

add_rect(s, Inches(6.85), Inches(1.5), Inches(5.9), Inches(5.35), DARK_BLUE, radius=0.05)
add_text(s, Inches(7.15), Inches(1.72), Inches(5.4), Inches(0.4), "Professional Skills", size=16, bold=True, color=SKY_BLUE)
prof_items = [
    "Problem solving under real constraints",
    "Structuring & organising a full project",
    "Following the software development lifecycle",
    "Reporting progress & issues honestly",
    "Time management against weekly deadlines",
]
tb = s.shapes.add_textbox(Inches(7.15), Inches(2.2), Inches(5.4), Inches(4.3))
tf = tb.text_frame; tf.word_wrap = True
for i, item in enumerate(prof_items):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(14)
    r = p.add_run(); r.text = "\u25A0  " + item
    r.font.size = Pt(13.5); r.font.color.rgb = WHITE; r.font.name = FONT

footer(s, 15)
add_notes(s, (
    "PRESENTER NOTES: This is a reflection slide — speak personally here rather than reading bullets. "
    "Give one concrete example per column (e.g. 'I learned to read a traceback bottom-up instead of "
    "guessing fixes').\n\n"
    "ANIMATION SUGGESTION: Left panel 'Fade', right panel 'Fly in from right'.\n"
    "TRANSITION SUGGESTION: Push (left)."
))

# =============================================================================
# SLIDE 16 — FUTURE ENHANCEMENTS
# =============================================================================
s = add_slide()
set_background(s, LIGHT_GRAY)
header_bar(s, "Roadmap", "Future Enhancements")

items = [
    ("\U0001F4B3", "Payment Gateway"), ("\U0001F531", "QR Ticket"), ("\U0001F4CD", "Live Bus Tracking"),
    ("\U0001F4F2", "OTP Login"), ("\U0001F4E7", "Email Notifications"), ("\u2B50", "Ratings & Reviews"),
    ("\U0001F4CA", "Admin Dashboard"), ("\U0001F4C8", "Booking Analytics"), ("\u2601", "Cloud Deployment"),
    ("\U0001F916", "AI Chatbot Support"), ("\U0001F510", "Password Hashing"), ("\U0001F5D1", "Online Cancellation"),
]
cols = 4
cw = Inches(2.95); ch = Inches(1.35); gx = Inches(0.15); gy = Inches(0.2)
x0 = Inches(0.55); y0 = Inches(1.6)
for i, (icon, name) in enumerate(items):
    col = i % cols
    row = i // cols
    x = x0 + col * (cw + gx)
    y = y0 + row * (ch + gy)
    add_rect(s, x, y, cw, ch, WHITE, line_color=MID_GRAY, radius=0.16)
    icon_badge(s, x + Inches(0.65), y + ch/2, Inches(0.65), icon, bg=ORANGE, icon_size=20)
    add_text(s, x + Inches(1.1), y, cw - Inches(1.25), ch, name, size=12.5, bold=True,
              color=DARK_BLUE, anchor=MSO_ANCHOR.MIDDLE)

footer(s, 16)
add_notes(s, (
    "PRESENTER NOTES: Present this as a realistic roadmap, not a wish-list — mention that these arise "
    "directly from limitations documented during testing (e.g. plain-text passwords \u2192 hashing).\n\n"
    "ANIMATION SUGGESTION: Grid tiles 'Wipe' in row by row, quick stagger.\n"
    "TRANSITION SUGGESTION: Fade."
))

# =============================================================================
# SLIDE 17 — PROJECT DEMONSTRATION
# =============================================================================
s = add_slide()
set_background(s, DARK_BLUE)
header_bar(s, "Live Walkthrough", "Project Demonstration")

add_text(s, Inches(0.55), Inches(1.5), Inches(12.2), Inches(0.5),
          "Live demo will show the complete booking journey on the running application.",
          size=15, color=SKY_BLUE)

demo_steps = ["Registration", "Login", "Search &\nBooking", "Payment", "Success\nPage", "My\nBookings"]
cw = Inches(1.95); ch = Inches(1.5); gx = Inches(0.2)
total_w = cw * 6 + gx * 5
x0 = (SLIDE_W - total_w) / 2
y0 = Inches(2.3)
for i, st in enumerate(demo_steps):
    x = x0 + i * (cw + gx)
    add_rect(s, x, y0, cw, ch, DARK_BLUE_2, line_color=SKY_BLUE, line_w=Pt(1), radius=0.15)
    icon_badge(s, x + cw/2, y0 + Inches(0.5), Inches(0.6), "\u25B6", bg=ORANGE, icon_size=18)
    add_text(s, x + Inches(0.05), y0 + Inches(0.95), cw - Inches(0.1), Inches(0.5), st,
              size=11.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_rect(s, Inches(0.55), Inches(4.3), Inches(12.25), Inches(2.55), DARK_BLUE_2, line_color=SKY_BLUE, radius=0.05)
add_text(s, Inches(0.85), Inches(4.5), Inches(11.6), Inches(0.4), "Talking Points During Demo", size=15, bold=True, color=ORANGE)
add_bullets(s, Inches(0.85), Inches(4.95), Inches(11.6), Inches(1.8), [
    "Point out the seat-count restriction enforced during seat selection",
    "Show the e-ticket generated immediately after payment confirmation",
    "Open My Bookings to prove the record was saved permanently in SQLite",
], size=13, color=WHITE, space_after=10)

footer(s, 17)
add_notes(s, (
    "PRESENTER NOTES: Switch to the live application (or a recorded screen capture as backup) here. "
    "Follow the six steps in order. Keep narration short — let the working software speak for "
    "itself. Have the DB Browser for SQLite open in a second window to show the row appearing live.\n\n"
    "ANIMATION SUGGESTION: None needed — this slide is a launch pad into the live demo, keep it "
    "static so you can switch windows without a distracting animation mid-flow.\n"
    "TRANSITION SUGGESTION: Cut / None (minimise distraction before switching to the live app)."
))

# =============================================================================
# SLIDE 18 — CONCLUSION
# =============================================================================
s = add_slide()
set_background(s, LIGHT_GRAY)
header_bar(s, "Wrap-Up", "Conclusion")

concl = [
    ("\U0001F393", "Internship Experience", "Hands-on exposure to a real development process,\nweekly reviews and mentor feedback"),
    ("\U0001F4A1", "Knowledge Gained", "Full-stack fundamentals: Flask, SQLite, HTML/CSS,\nJinja2, debugging & testing"),
    ("\U0001F3C1", "Project Outcome", "A fully working, end-to-end bus booking application\ndelivered and demonstrated"),
    ("\U0001F331", "Personal Growth", "Improved patience, discipline and honest\nself-assessment while debugging"),
    ("\U0001F4BC", "Career Impact", "A concrete, explainable project for interviews and\na clearer sense of web development as a career path"),
]
y = Inches(1.5)
for icon, title, desc in concl:
    add_rect(s, Inches(0.55), y, Inches(12.25), Inches(1.0), WHITE, line_color=MID_GRAY, radius=0.1)
    icon_badge(s, Inches(1.15), y + Inches(0.5), Inches(0.62), icon, bg=DARK_BLUE, icon_size=22)
    add_text(s, Inches(1.65), y + Inches(0.1), Inches(3.3), Inches(0.4), title, size=13.5, bold=True, color=DARK_BLUE)
    add_text(s, Inches(1.65), y + Inches(0.45), Inches(10.8), Inches(0.5), desc.replace("\n", "  "), size=10.5, color=TEXT_MUTED)
    y += Inches(1.1)

footer(s, 18)
add_notes(s, (
    "PRESENTER NOTES: Close the technical narrative here. Speak in first person and be genuine rather "
    "than reading verbatim — this slide should feel like a personal summary, not a repeat of earlier "
    "slides.\n\n"
    "ANIMATION SUGGESTION: Rows 'Fade' in one at a time, slower pace (0.3s) to let the panel absorb "
    "each point.\n"
    "TRANSITION SUGGESTION: Fade."
))

# =============================================================================
# SLIDE 19 — ACKNOWLEDGEMENT
# =============================================================================
s = add_slide()
set_background(s, DARK_BLUE)
header_bar(s, "Gratitude", "Acknowledgement")

ack = [
    ("\U0001F3E2", "DayLearner Private Limited", "For the opportunity to work on a real project"),
    ("\U0001F468\u200D\U0001F3EB", "Mentor", "For weekly guidance, code reviews & honest feedback"),
    ("\U0001F3EB", "GITAM University & Faculty", "For academic supervision and support throughout"),
    ("\U0001F465", "Friends", "For testing the application and reporting issues"),
    ("\u2764", "Parents", "For constant encouragement and patience"),
]
cols = 3
cw = Inches(3.95); ch = Inches(2.5); gx = Inches(0.2); gy = Inches(0.25)
x0 = Inches(0.6); y0 = Inches(1.6)
for i, (icon, title, desc) in enumerate(ack):
    col = i % cols
    row = i // cols
    x = x0 + col * (cw + gx)
    y = y0 + row * (ch + gy)
    add_rect(s, x, y, cw, ch, DARK_BLUE_2, line_color=SKY_BLUE, line_w=Pt(0.75), radius=0.08)
    icon_badge(s, x + cw/2, y + Inches(0.6), Inches(0.75), icon, bg=ORANGE, icon_size=26)
    add_text(s, x + Inches(0.15), y + Inches(1.15), cw - Inches(0.3), Inches(0.5), title,
              size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.2), y + Inches(1.65), cw - Inches(0.4), Inches(0.7), desc,
              size=10.5, color=SKY_BLUE, align=PP_ALIGN.CENTER)

footer(s, 19)
add_notes(s, (
    "PRESENTER NOTES: Read the names out if slide is being shown to the actual mentor/HOD/faculty in "
    "the room — fill in the mentor's actual name before presenting (currently a placeholder in the "
    "written report too).\n\n"
    "ANIMATION SUGGESTION: Cards 'Fly in from bottom' staggered.\n"
    "TRANSITION SUGGESTION: Fade."
))

# =============================================================================
# SLIDE 20 — THANK YOU
# =============================================================================
s = add_slide()
set_background(s, DARK_BLUE)
c1 = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.8), Inches(-1.5), Inches(4.5), Inches(4.5))
c1.fill.solid(); c1.fill.fore_color.rgb = DARK_BLUE_2; c1.line.fill.background(); c1.shadow.inherit = False
c2 = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-1.5), Inches(4.8), Inches(3.6), Inches(3.6))
c2.fill.solid(); c2.fill.fore_color.rgb = DARK_BLUE_2; c2.line.fill.background(); c2.shadow.inherit = False

icon_badge(s, SLIDE_W/2, Inches(2.15), Inches(1.0), "\U0001F68C", bg=ORANGE, icon_size=34)
add_text(s, 0, Inches(2.85), SLIDE_W, Inches(1.0), "Thank You", size=48, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, 0, Inches(3.7), SLIDE_W, Inches(0.5), "Questions & Discussion Welcome", size=18, color=SKY_BLUE, align=PP_ALIGN.CENTER)
add_rect(s, SLIDE_W/2 - Inches(0.8), Inches(4.35), Inches(1.6), Pt(3), ORANGE)

add_text(s, 0, Inches(4.7), SLIDE_W, Inches(0.4), "S. Sowmya  |  2023003449", size=14, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, 0, Inches(5.1), SLIDE_W, Inches(0.4), "GITAM (Deemed to be University), Hyderabad", size=12.5, color=SKY_BLUE, align=PP_ALIGN.CENTER)
add_text(s, 0, Inches(5.45), SLIDE_W, Inches(0.4), "Internship Project: EasyBus Pro  |  DayLearner Private Limited",
          size=12.5, color=SKY_BLUE, align=PP_ALIGN.CENTER)
add_text(s, 0, Inches(5.9), SLIDE_W, Inches(0.4), "<< student e-mail / contact placeholder >>",
          size=11.5, color=TEXT_MUTED, align=PP_ALIGN.CENTER, italic=True)

add_notes(s, (
    "PRESENTER NOTES: Stop talking once this slide is shown — pause, smile, and invite questions. "
    "Have the live application still open in a background tab in case the panel wants to see a "
    "specific screen again.\n\n"
    "ANIMATION SUGGESTION: Icon badge 'Zoom' in, title 'Fade' in after, subtitle lines follow with a "
    "gentle stagger.\n"
    "TRANSITION SUGGESTION: Fade to black (end of deck)."
))

# ---------------------------------------------------------------- save ----
out_path = r"c:\Users\reddy\OneDrive\Desktop\EasyBusPro\Internship_Report\EasyBus_Pro_Internship_Presentation.pptx"
prs.save(out_path)
print("Saved:", out_path)
print("Slide count:", len(prs.slides.__iter__.__self__._sldIdLst))
